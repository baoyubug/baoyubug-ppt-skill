#!/usr/bin/env python3
"""
模板应用脚本 - 将内容套用到PPT模板中
支持从JSON内容映射到模板占位符
"""

import os
import sys
import argparse
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


def apply_content_to_template(template_path, content_data, output_path):
    """
    将内容数据填充到PPT模板
    
    Args:
        template_path: 模板PPTX文件路径
        content_data: 内容数据（dict或json文件路径）
        output_path: 输出文件路径
    """
    
    # 加载模板
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")
    
    prs = Presentation(template_path)
    
    # 加载内容
    if isinstance(content_data, str) and os.path.exists(content_data):
        with open(content_data, 'r', encoding='utf-8') as f:
            content = json.load(f)
    elif isinstance(content_data, dict):
        content = content_data
    else:
        raise ValueError("content_data必须是dict或json文件路径")
    
    global_config = content.get("global", {})
    slides_content = content.get("slides", [])
    
    total_slides = len(prs.slides)
    
    # 逐页填充内容
    for i, slide_data in enumerate(slides_content):
        if i >= total_slides:
            print(f"警告: 内容页({len(slides_content)})多于模板页({total_slides})，跳过多余内容")
            break
        
        slide = prs.slides[i]
        fill_slide(slide, slide_data, global_config)
    
    # 如果内容页少于模板页，删除多余的模板页
    if len(slides_content) < total_slides:
        for i in range(total_slides - 1, len(slides_content) - 1, -1):
            slide_id = prs.slides._sldIdLst[i]
            prs.part.drop_rel(slide_id.rId)
            del prs.slides._sldIdLst[i]
        print(f"已删除 {total_slides - len(slides_content)} 页多余的模板幻灯片")
    
    # 保存
    prs.save(output_path)
    print(f"模板填充完成，已保存到: {output_path}")
    return output_path


def fill_slide(slide, slide_data, global_config=None):
    """
    填充单页幻灯片内容
    
    支持三种填充方式：
    1. 按placeholder index匹配
    2. 按形状名称匹配
    3. 按位置顺序填充
    """
    if global_config is None:
        global_config = {}
    
    # 获取填充映射
    mapping = slide_data.get("mapping", {})
    placeholder_map = slide_data.get("placeholders", {})
    
    # 方式1: 按placeholder index匹配
    if placeholder_map:
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                idx = str(shape.placeholder_format.idx)
                if idx in placeholder_map:
                    fill_shape_text(shape, placeholder_map[idx], global_config)
    
    # 方式2: 按形状名称匹配
    if mapping:
        for shape in slide.shapes:
            shape_name = shape.name
            if shape_name in mapping:
                fill_shape_text(shape, mapping[shape_name], global_config)
    
    # 方式3: 如果没指定映射，按顺序填充文本占位符
    if not placeholder_map and not mapping and "content" in slide_data:
        content_list = slide_data["content"]
        if isinstance(content_list, str):
            content_list = [content_list]
        
        text_index = 0
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                if text_index < len(content_list):
                    fill_shape_text(shape, content_list[text_index], global_config)
                    text_index += 1


def fill_shape_text(shape, text_content, global_config):
    """填充形状的文本内容"""
    if not hasattr(shape, 'text_frame'):
        return
    
    # 如果text_content是dict，包含text和style
    if isinstance(text_content, dict):
        text = text_content.get("text", "")
        style = text_content.get("style", {})
    else:
        text = str(text_content)
        style = {}
    
    text_frame = shape.text_frame
    
    # 如果只有一段文本，直接设置
    if '\n' not in text and not text_frame.paragraphs[0].runs:
        text_frame.paragraphs[0].text = ""
        run = text_frame.paragraphs[0].add_run()
        run.text = text
        apply_run_style(run, style, global_config)
    else:
        # 多段文本
        text_frame.clear()
        lines = text.split('\n')
        for line in lines:
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = line
            apply_run_style(run, style, global_config)
        
        # 删除第一个空段落（clear后默认段落）
        if len(text_frame.paragraphs) > len(lines):
            trPr = text_frame.paragraphs[0]._p
            trPr.getparent().remove(trPr)


def apply_run_style(run, style, global_config):
    """应用文本样式"""
    # 优先使用局部样式，回退到全局配置
    font_size = style.get("font_size", global_config.get("font_size"))
    font_color = style.get("color", global_config.get("color"))
    font_bold = style.get("bold", global_config.get("bold"))
    
    if font_size:
        run.font.size = Pt(font_size)
    if font_color:
        run.font.color.rgb = RGBColor(*font_color)
    if font_bold is not None:
        run.font.bold = font_bold


def main():
    parser = argparse.ArgumentParser(description="将内容套用到PPT模板")
    parser.add_argument("--template", "-t", required=True, help="模板PPTX文件路径")
    parser.add_argument("--content", "-c", required=True, help="内容JSON文件路径")
    parser.add_argument("--output", "-o", required=True, help="输出PPTX文件路径")
    
    args = parser.parse_args()
    
    try:
        result = apply_content_to_template(args.template, args.content, args.output)
        print(f"成功: {result}")
        return 0
    except Exception as e:
        print(f"模板应用失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == "__main__":
    sys.exit(main())
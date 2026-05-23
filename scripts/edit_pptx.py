#!/usr/bin/env python3
"""
PPT编辑脚本 - 对已有.pptx文件进行编辑操作
支持：替换文本、调整布局、增删页面、修改样式
"""

import os
import sys
import argparse
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 确保中文输出无乱码（Windows PowerShell）
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')


def load_presentation(file_path):
    """加载PPT文件"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    return Presentation(file_path)


def edit_slide_text(prs, slide_index, edits):
    """
    编辑指定幻灯片的文本
    
    Args:
        prs: Presentation对象
        slide_index: 幻灯片索引（0-based）
        edits: 编辑操作列表 [{"old_text": "...", "new_text": "...", "placeholder_index": 0}, ...]
    """
    if slide_index >= len(prs.slides):
        raise IndexError(f"幻灯片索引超出范围: {slide_index} (共{len(prs.slides)}页)")
    
    slide = prs.slides[slide_index]
    
    for edit in edits:
        old_text = edit.get("old_text", "")
        new_text = edit.get("new_text", "")
        shape_index = edit.get("shape_index", None)
        
        if shape_index is not None and shape_index < len(slide.shapes):
            shape = slide.shapes[shape_index]
            if hasattr(shape, 'text') and old_text in shape.text:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                elif shape.has_text_frame:
                    shape.text = shape.text.replace(old_text, new_text)
        else:
            # 遍历所有shapes查找
            for shape in slide.shapes:
                if hasattr(shape, 'text') and old_text in shape.text:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    break
                    elif shape.has_text_frame:
                        shape.text = shape.text.replace(old_text, new_text)
    
    print(f"已完成第{slide_index+1}页的文本编辑")


def delete_slide(prs, slide_index):
    """删除指定幻灯片"""
    if slide_index >= len(prs.slides):
        raise IndexError(f"幻灯片索引超出范围: {slide_index}")
    
    slide_id = prs.slides._sldIdLst[slide_index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[slide_index]
    
    print(f"已删除第{slide_index+1}页幻灯片")


def reorder_slides(prs, from_index, to_index):
    """调整幻灯片顺序"""
    if from_index >= len(prs.slides) or to_index >= len(prs.slides):
        raise IndexError("幻灯片索引超出范围")
    
    sld_list = list(prs.slides._sldIdLst)
    item = sld_list.pop(from_index)
    sld_list.insert(to_index, item)
    prs.slides._sldIdLst.clear()
    for sld in sld_list:
        prs.slides._sldIdLst.append(sld)
    
    print(f"已将第{from_index+1}页移动到第{to_index+1}页位置")


def duplicate_slide(prs, slide_index):
    """复制幻灯片"""
    if slide_index >= len(prs.slides):
        raise IndexError("幻灯片索引超出范围")
    
    slide = prs.slides[slide_index]
    slide_layout = slide.slide_layout
    
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 复制内容
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.has_text_frame:
            # 找到新幻灯片中对应的占位符
            for new_shape in new_slide.shapes:
                if hasattr(new_shape, 'placeholder_format') and hasattr(shape, 'placeholder_format'):
                    if new_shape.placeholder_format.idx == shape.placeholder_format.idx:
                        new_shape.text = shape.text
                        break
    
    print(f"已复制第{slide_index+1}页幻灯片")


def apply_style(prs, slide_index, style_config):
    """
    应用样式到指定幻灯片
    
    style_config格式: {"font_size": 20, "color": [255,105,180], "alignment": "center", ...}
    """
    if slide_index >= len(prs.slides):
        raise IndexError("幻灯片索引超出范围")
    
    slide = prs.slides[slide_index]
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "font_size" in style_config:
                        run.font.size = Pt(style_config["font_size"])
                    if "color" in style_config:
                        run.font.color.rgb = RGBColor(*style_config["color"])
                    if "bold" in style_config:
                        run.font.bold = style_config["bold"]
                if "alignment" in style_config:
                    align_map = {
                        "left": PP_ALIGN.LEFT,
                        "center": PP_ALIGN.CENTER,
                        "right": PP_ALIGN.RIGHT,
                        "justify": PP_ALIGN.JUSTIFY
                    }
                    paragraph.alignment = align_map.get(style_config["alignment"], PP_ALIGN.LEFT)
    
    print(f"已对第{slide_index+1}页应用样式")


def main():
    parser = argparse.ArgumentParser(description="编辑PPT文件")
    parser.add_argument("--file", "-f", required=True, help="PPTX文件路径")
    parser.add_argument("--action", "-a", required=True, 
                        choices=["edit_text", "delete_slide", "reorder", "duplicate", "apply_style"])
    parser.add_argument("--params", "-p", required=True, help="操作参数JSON字符串")
    parser.add_argument("--output", "-o", required=True, help="输出文件路径")
    
    args = parser.parse_args()
    
    try:
        prs = load_presentation(args.file)
        params = json.loads(args.params)
        
        if args.action == "edit_text":
            slide_index = params.get("slide_index", 0)
            edits = params.get("edits", [])
            edit_slide_text(prs, slide_index, edits)
        
        elif args.action == "delete_slide":
            slide_index = params.get("slide_index", 0)
            delete_slide(prs, slide_index)
        
        elif args.action == "reorder":
            from_index = params.get("from_index", 0)
            to_index = params.get("to_index", 0)
            reorder_slides(prs, from_index, to_index)
        
        elif args.action == "duplicate":
            slide_index = params.get("slide_index", 0)
            duplicate_slide(prs, slide_index)
        
        elif args.action == "apply_style":
            slide_index = params.get("slide_index", 0)
            style_config = params.get("style_config", {})
            apply_style(prs, slide_index, style_config)
        
        prs.save(args.output)
        print(f"编辑完成，已保存到: {args.output}")
        return 0
    
    except Exception as e:
        print(f"编辑失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == "__main__":
    sys.exit(main())
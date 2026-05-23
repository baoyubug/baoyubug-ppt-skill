#!/usr/bin/env python3
"""
PPT生成脚本 - 小红书清新风格
支持从主题生成PPT、自动排版、应用模板
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import argparse

# 确保中文输出无乱码（Windows PowerShell）
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

# 小红书清新风配色方案
XHS_COLORS = {
    "primary": RGBColor(255, 105, 180),    # 粉红
    "secondary": RGBColor(100, 200, 200),   # 青绿  
    "accent": RGBColor(255, 215, 0),       # 金色
    "text": RGBColor(50, 50, 50),          # 深灰
    "background": RGBColor(255, 250, 245), # 米白
    "light_bg": RGBColor(255, 253, 250),   # 浅米白
    "border": RGBColor(230, 230, 230)      # 浅灰边框
}

def create_xhs_style_presentation(output_path, content_data, template_path=None):
    """
    创建小红书清新风格PPT
    
    Args:
        output_path: 输出文件路径
        content_data: 内容数据，可以是dict或json文件路径
        template_path: 可选模板路径
    """
    
    # 加载内容
    if isinstance(content_data, str) and os.path.exists(content_data):
        with open(content_data, 'r', encoding='utf-8') as f:
            content = json.load(f)
    elif isinstance(content_data, dict):
        content = content_data
    else:
        raise ValueError("content_data必须是dict或json文件路径")
    
    # 创建或加载PPT
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # 设置默认尺寸为16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    # 提取内容
    title = content.get("title", "未命名演示文稿")
    slides_content = content.get("slides", [])
    theme = content.get("theme", "xhs_fresh")  # 小红书清新风
    
    print(f"正在生成PPT: {title}")
    print(f"包含 {len(slides_content)} 页幻灯片")
    
    # 创建封面页
    create_cover_slide(prs, title, content.get("subtitle", ""))
    
    # 创建目录页（如果有多个章节）
    if len(slides_content) > 3:
        create_toc_slide(prs, slides_content)
    
    # 创建内容页
    for i, slide_data in enumerate(slides_content):
        create_content_slide(prs, slide_data, i+1)
    
    # 创建结束页
    create_ending_slide(prs, content.get("ending_text", "感谢观看"))
    
    # 保存文件
    prs.save(output_path)
    print(f"PPT已保存到: {output_path}")
    return output_path

def create_cover_slide(prs, title, subtitle=""):
    """创建封面页"""
    # 使用标题布局
    slide_layout = prs.slide_layouts[0]  # 标题布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = XHS_COLORS["text"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 设置副标题
    if subtitle and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = XHS_COLORS["secondary"]
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加装饰元素
    add_decorative_elements(slide)

def create_toc_slide(prs, slides_content):
    """创建目录页"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = "目录"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = XHS_COLORS["primary"]
    
    # 添加目录内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()  # 清空默认内容
    
    for i, slide_data in enumerate(slides_content[:8]):  # 最多显示8个
        slide_title = slide_data.get("title", f"第{i+1}页")
        p = text_frame.add_paragraph()
        p.text = f"{i+1}. {slide_title}"
        p.font.size = Pt(20)
        p.font.color.rgb = XHS_COLORS["text"]
        p.level = 0

def create_content_slide(prs, slide_data, slide_num):
    """创建内容页"""
    layout_type = slide_data.get("layout", "title_content")
    
    if layout_type == "title_only":
        slide_layout = prs.slide_layouts[5]  # 仅标题
    elif layout_type == "two_content":
        slide_layout = prs.slide_layouts[3]  # 两栏内容
    else:
        slide_layout = prs.slide_layouts[1]  # 标题和内容
    
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = slide_data.get("title", f"第{slide_num}页")
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = XHS_COLORS["primary"]
    
    # 根据布局类型设置内容
    if layout_type == "title_content":
        if "content" in slide_data and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            set_content_text(content_shape, slide_data["content"])
    
    elif layout_type == "two_content":
        if "left_content" in slide_data and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            left_shape = slide.placeholders[1]
            set_content_text(left_shape, slide_data["left_content"])
        
        if "right_content" in slide_data and hasattr(slide, 'placeholders') and len(slide.placeholders) > 2:
            right_shape = slide.placeholders[2]
            set_content_text(right_shape, slide_data["right_content"])
    
    # 添加页码
    add_page_number(slide, slide_num)

def set_content_text(shape, content):
    """设置内容文本"""
    if isinstance(content, list):
        text = "\n".join(content)
    else:
        text = str(content)
    
    shape.text = text
    text_frame = shape.text_frame
    
    # 设置文本样式
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = XHS_COLORS["text"]
        paragraph.space_after = Pt(12)

def create_ending_slide(prs, ending_text):
    """创建结束页"""
    slide_layout = prs.slide_layouts[0]  # 标题布局
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = ending_text
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = XHS_COLORS["accent"]
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_decorative_elements(slide):
    """添加装饰元素"""
    # 在底部添加装饰线
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(11.33)
    height = Inches(0.05)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = XHS_COLORS["border"]
    line.line.fill.background()

def add_page_number(slide, page_num):
    """添加页码"""
    left = Inches(12)
    top = Inches(6.8)
    width = Inches(1)
    height = Inches(0.4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.color.rgb = XHS_COLORS["border"]
    p.alignment = PP_ALIGN.RIGHT

def main():
    parser = argparse.ArgumentParser(description="生成小红书清新风格PPT")
    parser.add_argument("--output", "-o", required=True, help="输出PPT文件路径")
    parser.add_argument("--content", "-c", required=True, help="内容JSON文件路径")
    parser.add_argument("--template", "-t", help="模板PPT文件路径（可选）")
    
    args = parser.parse_args()
    
    try:
        result = create_xhs_style_presentation(
            args.output,
            args.content,
            args.template
        )
        print(f"生成成功: {result}")
        return 0
    except Exception as e:
        print(f"生成失败: {str(e)}")
        return 1

if __name__ == "__main__":
    sys.exit(main())